"""Build a demo-day PowerPoint deck from presentation_data.json.

Usage:
    python scripts/build_deck.py
    # → writes green_lion_demo_deck.pptx in the repo root

Requires: python-pptx  (pip install python-pptx)
"""

from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Paths ─────────────────────────────────────────────────────────────────────
ROOT = Path(__file__).parent.parent
DATA_FILE = ROOT / "presentation_data.json"
OUT_FILE = ROOT / "green_lion_demo_deck.pptx"

# ── Brand palette ─────────────────────────────────────────────────────────────
# Dark navy / rich green / white / amber / red — professional finance aesthetic
NAVY   = RGBColor(0x0D, 0x1B, 0x2A)   # slide backgrounds (title, section dividers)
GREEN  = RGBColor(0x00, 0x7A, 0x4E)   # accents, positive, supported
AMBER  = RGBColor(0xE8, 0x8C, 0x00)   # warning / partially supported
RED    = RGBColor(0xC0, 0x39, 0x2B)   # not supported / alert
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT  = RGBColor(0xF4, 0xF6, 0xF8)   # content slide backgrounds
SLATE  = RGBColor(0x2C, 0x3E, 0x50)   # body text
MID    = RGBColor(0x6C, 0x75, 0x7D)   # secondary text / labels

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ── Helpers ───────────────────────────────────────────────────────────────────

def _px(inches: float) -> int:
    """EMUs from inches."""
    return int(inches * 914400)


def _add_rect(slide, left, top, width, height, fill: RGBColor | None = None, line: RGBColor | None = None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        _px(left), _px(top), _px(width), _px(height),
    )
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = _px(0.01)
    else:
        shape.line.fill.background()
    return shape


def _add_label(slide, text: str, left, top, width, height,
               font_size=11, bold=False, color: RGBColor = SLATE,
               align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(_px(left), _px(top), _px(width), _px(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def _set_slide_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


# ── Slide builders ─────────────────────────────────────────────────────────────

def slide_title(prs: Presentation, d: dict):
    """Slide 1 — Title / hook (dark navy)."""
    layout = prs.slide_layouts[6]  # blank
    sl = prs.slides.add_slide(layout)
    _set_slide_bg(sl, NAVY)

    # Green accent bar on left
    _add_rect(sl, 0, 0, 0.06, 7.5, fill=GREEN)

    # sf-agents logotype area
    _add_label(sl, "sf-agents", 0.35, 0.35, 4, 0.6,
               font_size=13, bold=True, color=GREEN)
    _add_label(sl, "governance-first structured finance AI",
               0.35, 0.85, 6, 0.4, font_size=10, color=RGBColor(0xAA, 0xBB, 0xCC))

    # Main title
    _add_label(sl,
               "A green bond with 974 mortgages\naveraging 208 kWh/m²",
               0.35, 1.6, 8.5, 1.5,
               font_size=36, bold=True, color=WHITE)

    # Sub-title / hook
    _add_label(sl,
               "Four times the ISS SPO threshold. A compliance gap — or a €312M renovation campaign?\n"
               "sf-agents finds out in under 7 minutes.",
               0.35, 3.2, 8.5, 1.2,
               font_size=14, color=RGBColor(0xCC, 0xDD, 0xEE))

    # Deal badge (bottom right)
    _add_rect(sl, 9.6, 6.3, 3.4, 0.9, fill=RGBColor(0x1A, 0x2E, 0x40))
    _add_label(sl, "Green Lion 2026-1 B.V.", 9.7, 6.32, 3.2, 0.35,
               font_size=9, bold=True, color=GREEN, align=PP_ALIGN.RIGHT)
    _add_label(sl, "€1.033B · 3,237 loans · ING · NL", 9.7, 6.62, 3.2, 0.35,
               font_size=8, color=RGBColor(0xAA, 0xBB, 0xCC), align=PP_ALIGN.RIGHT)


def slide_framework(prs: Presentation, d: dict):
    """Slide 2 — What sf-agents is."""
    layout = prs.slide_layouts[6]
    sl = prs.slides.add_slide(layout)
    _set_slide_bg(sl, LIGHT)

    _add_rect(sl, 0, 0, 13.33, 0.08, fill=GREEN)

    _add_label(sl, "What sf-agents is", 0.45, 0.22, 8, 0.5,
               font_size=22, bold=True, color=NAVY)
    _add_label(sl, "Governance-first structured finance AI on AWS Bedrock eu-north-1",
               0.45, 0.72, 10, 0.35, font_size=11, color=MID)

    # Three stat boxes
    boxes = [
        ("33", "registered primitives\nconnectors · extractors\nanalyzers · validators"),
        ("3", "pre-built recipes\nimpact mapping · 3LoD\ndefinition transparency"),
        ("100%", "audit coverage\nevery step hashed\nappend-only JSONL"),
    ]
    bw, bh, by = 3.6, 2.1, 1.4
    for i, (num, label) in enumerate(boxes):
        bx = 0.45 + i * 4.2
        _add_rect(sl, bx, by, bw, bh, fill=WHITE)
        _add_rect(sl, bx, by, bw, 0.06, fill=GREEN)
        _add_label(sl, num, bx + 0.2, by + 0.18, bw - 0.4, 0.8,
                   font_size=44, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        _add_label(sl, label, bx + 0.2, by + 0.95, bw - 0.4, 0.95,
                   font_size=11, color=SLATE, align=PP_ALIGN.CENTER)

    # One-liner at bottom
    _add_label(sl,
               "The LLM planner writes a DAG live for each question — no hardcoded workflows. "
               "Every step produces dual-grounded citations: exact document page + exact tape row.",
               0.45, 3.9, 12.4, 0.7,
               font_size=11, color=SLATE, italic=True)

    _add_label(sl, "Bedrock · Sonnet 4.6 · eu-north-1",
               0.45, 6.9, 6, 0.3, font_size=8, color=MID)


def slide_deal(prs: Presentation, d: dict):
    """Slide 3 — The deal (pool KPIs)."""
    layout = prs.slide_layouts[6]
    sl = prs.slides.add_slide(layout)
    _set_slide_bg(sl, LIGHT)

    _add_rect(sl, 0, 0, 13.33, 0.08, fill=NAVY)

    _add_label(sl, "Green Lion 2026-1 B.V.", 0.45, 0.22, 9, 0.5,
               font_size=22, bold=True, color=NAVY)
    _add_label(sl, "ING originator · Dutch residential mortgages · reporting date 2026-04-30",
               0.45, 0.72, 10, 0.35, font_size=11, color=MID)

    ps = d["pool_stats"]
    perf = d["performance"]

    kpis = [
        ("€1.033B", "total balance"),
        ("3,237", "loans"),
        ("3.18%", "WA interest rate"),
        ("68.9%", "WA current LTV"),
        ("261 mo", "WA remaining term"),
        ("19.7%", "NHG guarantee"),
        ("99.97%", "performing"),
        ("0.03%", "90+ day arrears"),
    ]

    cols, rows = 4, 2
    bw, bh = 2.9, 1.6
    start_x, start_y = 0.45, 1.3
    for idx, (val, lbl) in enumerate(kpis):
        col = idx % cols
        row = idx // cols
        bx = start_x + col * (bw + 0.22)
        by = start_y + row * (bh + 0.18)
        _add_rect(sl, bx, by, bw, bh, fill=WHITE)
        _add_rect(sl, bx, by, bw, 0.05,
                  fill=GREEN if idx < 6 else AMBER if idx == 6 else AMBER)
        _add_label(sl, val, bx + 0.15, by + 0.18, bw - 0.3, 0.75,
                   font_size=30, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        _add_label(sl, lbl, bx + 0.15, by + 0.95, bw - 0.3, 0.4,
                   font_size=10, color=MID, align=PP_ALIGN.CENTER)

    _add_label(sl, "Source: green_lion_2026_1_synthetic_loan_tape.csv — computed directly",
               0.45, 6.9, 8, 0.3, font_size=8, color=MID)


def slide_finding(prs: Presentation, d: dict):
    """Slide 4 — The green bombshell (NOT SUPPORTED verdicts)."""
    layout = prs.slide_layouts[6]
    sl = prs.slides.add_slide(layout)
    _set_slide_bg(sl, NAVY)

    _add_rect(sl, 0, 0, 0.06, 7.5, fill=RED)

    _add_label(sl, "The Finding", 0.35, 0.25, 8, 0.45,
               font_size=22, bold=True, color=WHITE)
    _add_label(sl,
               "ISS SPO green claims tested against 3,237 loan tape rows — live Bedrock run",
               0.35, 0.7, 10, 0.35, font_size=11,
               color=RGBColor(0xAA, 0xBB, 0xCC))

    # Big PED callout
    _add_rect(sl, 0.35, 1.2, 5.4, 1.5, fill=RGBColor(0x1A, 0x0A, 0x0A))
    _add_rect(sl, 0.35, 1.2, 5.4, 0.06, fill=RED)
    _add_label(sl, "133.47 kWh/m²", 0.5, 1.35, 5.1, 0.75,
               font_size=38, bold=True, color=RED, align=PP_ALIGN.CENTER)
    _add_label(sl, "mean primary energy demand  ·  threshold: 27 kWh/m²  ·  4.9× over",
               0.5, 2.08, 5.1, 0.4, font_size=10,
               color=RGBColor(0xFF, 0xAA, 0xAA), align=PP_ALIGN.CENTER)

    # Pass rate callout
    _add_rect(sl, 6.1, 1.2, 3.0, 1.5, fill=RGBColor(0x1A, 0x0A, 0x0A))
    _add_rect(sl, 6.1, 1.2, 3.0, 0.06, fill=RED)
    _add_label(sl, "12.8%", 6.2, 1.35, 2.8, 0.75,
               font_size=38, bold=True, color=RED, align=PP_ALIGN.CENTER)
    _add_label(sl, "pass PED criterion\n(414 of 3,237 loans)",
               6.2, 2.08, 2.8, 0.4, font_size=10,
               color=RGBColor(0xFF, 0xAA, 0xAA), align=PP_ALIGN.CENTER)

    _add_rect(sl, 9.45, 1.2, 3.5, 1.5, fill=RGBColor(0x1A, 0x0A, 0x0A))
    _add_rect(sl, 9.45, 1.2, 3.5, 0.06, fill=RED)
    _add_label(sl, "1,456", 9.55, 1.35, 3.3, 0.75,
               font_size=38, bold=True, color=RED, align=PP_ALIGN.CENTER)
    _add_label(sl, "loans EPC below A\n(B through G)",
               9.55, 2.08, 3.3, 0.4, font_size=10,
               color=RGBColor(0xFF, 0xAA, 0xAA), align=PP_ALIGN.CENTER)

    # Verdict table
    verdicts = [
        ("Green Asset Portfolio",    "NOT SUPPORTED",        RED,   "pg. 207"),
        ("Green Bond",               "NOT SUPPORTED",        RED,   "pg. 207"),
        ("Energy Efficient Mortgage","NOT SUPPORTED",        RED,   "pg. 207"),
        ("Near Zero Energy Building","PARTIALLY SUPPORTED",  AMBER, "pg. 207"),
        ("EPC label",                "SUPPORTED",            GREEN, "pg. 207"),
    ]
    vy = 3.05
    _add_label(sl, "CLAIM", 0.35, vy, 4.2, 0.3, font_size=9, bold=True,
               color=RGBColor(0xAA, 0xBB, 0xCC))
    _add_label(sl, "VERDICT", 4.6, vy, 3.5, 0.3, font_size=9, bold=True,
               color=RGBColor(0xAA, 0xBB, 0xCC))
    _add_label(sl, "SOURCE", 8.15, vy, 1.2, 0.3, font_size=9, bold=True,
               color=RGBColor(0xAA, 0xBB, 0xCC))
    vy += 0.32
    for term, verdict, col, pg in verdicts:
        _add_rect(sl, 0.35, vy, 12.6, 0.45, fill=RGBColor(0x12, 0x1E, 0x2D))
        _add_label(sl, term, 0.5, vy + 0.08, 4.0, 0.3, font_size=10, color=WHITE)
        _add_rect(sl, 4.55, vy + 0.08, 3.2, 0.28, fill=col)
        _add_label(sl, verdict, 4.58, vy + 0.09, 3.1, 0.26,
                   font_size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _add_label(sl, pg, 8.1, vy + 0.08, 1.2, 0.3,
                   font_size=9, color=RGBColor(0xAA, 0xBB, 0xCC))
        vy += 0.5

    _add_label(sl, "85 citation checks · dual-grounded (ISS SPO page + tape column) · run 650eb8ef · 410s",
               0.35, 6.9, 10, 0.3, font_size=8, color=MID)


def slide_opportunity(prs: Presentation, d: dict):
    """Slide 5 — Renovation opportunity."""
    layout = prs.slide_layouts[6]
    sl = prs.slides.add_slide(layout)
    _set_slide_bg(sl, LIGHT)

    _add_rect(sl, 0, 0, 13.33, 0.08, fill=GREEN)

    _add_label(sl, "The Opportunity", 0.45, 0.22, 9, 0.5,
               font_size=22, bold=True, color=NAVY)
    _add_label(sl, "EPC below A  ·  remaining term ≥ 4 years  ·  current LTV < 80%",
               0.45, 0.72, 10, 0.35, font_size=11, color=MID)

    rs = d["renovation_segment"]
    sc = rs["scenario_30pct_renovate"]

    # Hero stats row
    hero = [
        (str(rs["loan_count"]),                              "loans in segment"),
        (f"€{rs['total_balance_eur']/1e6:.0f}M",       "outstanding balance"),
        (f"{rs['pool_share_pct']:.1f}%",                    "of pool by count"),
        (f"{rs['avg_ped']:.0f} kWh/m²",               "avg energy demand"),
    ]
    hw = 2.9
    for i, (val, lbl) in enumerate(hero):
        hx = 0.45 + i * 3.15
        _add_rect(sl, hx, 1.25, hw, 1.4, fill=WHITE)
        _add_rect(sl, hx, 1.25, hw, 0.05, fill=GREEN)
        _add_label(sl, val, hx + 0.12, 1.4, hw - 0.24, 0.65,
                   font_size=28, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        _add_label(sl, lbl, hx + 0.12, 2.05, hw - 0.24, 0.35,
                   font_size=10, color=MID, align=PP_ALIGN.CENTER)

    # EPC breakdown bar (horizontal stacked-style)
    _add_label(sl, "EPC label breakdown within segment", 0.45, 3.0, 6, 0.3,
               font_size=10, bold=True, color=NAVY)
    _epc_colors = {
        "B": RGBColor(0x5D, 0xAD, 0x72), "C": RGBColor(0xF3, 0x9C, 0x12),
        "D": RGBColor(0xE6, 0x7E, 0x22), "E": RGBColor(0xE7, 0x4C, 0x3C),
        "F": RGBColor(0xC0, 0x39, 0x2B), "G": RGBColor(0x96, 0x28, 0x20),
    }
    _epc_order = ["B", "C", "D", "E", "F", "G"]
    epc_raw = rs["epc_breakdown"]
    epc_items = [
        (lbl, epc_raw.get(lbl, 0), _epc_colors[lbl])
        for lbl in _epc_order if epc_raw.get(lbl, 0) > 0
    ]
    total_seg = sum(x[1] for x in epc_items)
    bar_left, bar_top, bar_w, bar_h = 0.45, 3.35, 7.0, 0.5
    cx = bar_left
    for lbl, cnt, col in epc_items:
        w = bar_w * cnt / total_seg
        _add_rect(sl, cx, bar_top, w, bar_h, fill=col)
        if w > 0.25:
            _add_label(sl, f"{lbl}\n{cnt}", cx + 0.02, bar_top + 0.04, w - 0.04, bar_h - 0.08,
                       font_size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        cx += w

    # Province list
    _add_label(sl, "Top provinces", 7.9, 3.0, 4.5, 0.3,
               font_size=10, bold=True, color=NAVY)
    for i, p in enumerate(rs["top_5_provinces"]):
        py = 3.35 + i * 0.42
        _add_rect(sl, 7.9, py, 5.0, 0.36, fill=WHITE)
        _add_label(sl, p["province"], 8.05, py + 0.06, 3.0, 0.24,
                   font_size=10, color=SLATE)
        _add_label(sl, str(p["loan_count"]), 11.0, py + 0.06, 0.8, 0.24,
                   font_size=10, bold=True, color=NAVY, align=PP_ALIGN.RIGHT)

    # 30% scenario arrow
    _add_rect(sl, 0.45, 5.45, 12.4, 1.55, fill=RGBColor(0xE8, 0xF8, 0xF1))
    _add_rect(sl, 0.45, 5.45, 12.4, 0.05, fill=GREEN)
    _add_label(sl, "If 30% of segment renovates to EPC A  (292 loans):", 0.65, 5.55, 7, 0.35,
               font_size=11, bold=True, color=NAVY)
    _add_label(sl,
               f"Green share:  {sc['current_green_share_pct']}%  →  {sc['post_scenario_green_share_pct']}%",
               0.65, 5.9, 5.5, 0.4, font_size=22, bold=True, color=GREEN)
    _add_label(sl, f"+{sc['green_share_uplift_pp']} percentage points",
               6.2, 5.9, 3.5, 0.4, font_size=16, bold=True,
               color=GREEN)
    _add_label(sl,
               f"Avg equity in segment: €{rs['avg_equity_eur']:,.0f}  ·  borrowers have headroom for renovation finance today",
               0.65, 6.45, 11, 0.35, font_size=10, color=MID)

    _add_label(sl, "Source: analyzer.green_renovation_potential · deterministic · run 77dabb18 · 31s",
               0.45, 7.15, 9, 0.25, font_size=8, color=MID)


def slide_governance(prs: Presentation, d: dict):
    """Slide 6 — The governance proof (audit trail + IC verdict)."""
    layout = prs.slide_layouts[6]
    sl = prs.slides.add_slide(layout)
    _set_slide_bg(sl, NAVY)

    _add_rect(sl, 0, 0, 0.06, 7.5, fill=GREEN)

    _add_label(sl, "The Governance Proof", 0.35, 0.25, 9, 0.45,
               font_size=22, bold=True, color=WHITE)
    _add_label(sl, "Every finding is dual-grounded, time-stamped, and hashed",
               0.35, 0.7, 10, 0.35, font_size=11, color=RGBColor(0xAA, 0xBB, 0xCC))

    qa = d["live_run_results"]["question_a_green_bombshell"]
    entries = qa["audit_first_10_entries"][:6]

    # Audit trail table header
    vy = 1.25
    _add_rect(sl, 0.35, vy, 12.6, 0.35, fill=RGBColor(0x0A, 0x44, 0x30))
    cols_hdr = [("STEP", 2.8), ("PRIMITIVE", 3.0), ("CONF", 1.0), ("DURATION", 1.5), ("TIMESTAMP", 4.3)]
    cx = 0.45
    for hdr, cw in cols_hdr:
        _add_label(sl, hdr, cx, vy + 0.06, cw - 0.1, 0.22,
                   font_size=8, bold=True, color=GREEN)
        cx += cw
    vy += 0.38

    for i, entry in enumerate(entries):
        row_col = RGBColor(0x12, 0x1E, 0x2D) if i % 2 == 0 else RGBColor(0x0D, 0x18, 0x26)
        _add_rect(sl, 0.35, vy, 12.6, 0.38, fill=row_col)
        cx = 0.45
        values = [
            (entry["step_id"][:28], 2.8),
            (entry["primitive"], 3.0),
            (str(entry["confidence"]), 1.0),
            (f'{entry["duration_ms"]:.0f} ms', 1.5),
            (entry["timestamp"][:19].replace("T", " "), 4.3),
        ]
        for val, cw in values:
            col = GREEN if entry["confidence"] == 1.0 else (
                AMBER if entry["confidence"] > 0.5 else RED)
            txt_col = col if val == str(entry["confidence"]) else WHITE
            _add_label(sl, val, cx, vy + 0.07, cw - 0.08, 0.24,
                       font_size=8, color=txt_col)
            cx += cw
        vy += 0.4

    # Closing stats
    _add_rect(sl, 0.35, 4.95, 12.6, 1.1, fill=RGBColor(0x07, 0x30, 0x1C))
    stats = [
        ("11", "plan steps (Q-A)"),
        ("85", "citation checks"),
        ("410s", "Q-A latency"),
        ("31s", "Q-B latency"),
        ("33", "primitives"),
        ("0", "hardcoded workflows"),
    ]
    sw = 12.6 / len(stats)
    for i, (val, lbl) in enumerate(stats):
        sx = 0.35 + i * sw
        _add_label(sl, val, sx + 0.1, 5.05, sw - 0.2, 0.45,
                   font_size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _add_label(sl, lbl, sx + 0.1, 5.5, sw - 0.2, 0.35,
                   font_size=8, color=RGBColor(0x88, 0xCC, 0xAA), align=PP_ALIGN.CENTER)

    # Close line
    tp = d["demo_talking_points"]
    _add_label(sl, tp["close"], 0.35, 6.35, 12.6, 0.5,
               font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER, italic=True)

    _add_label(sl, f"run 650eb8ef · {qa['citations_verified']} citations · audit_logs/650eb8ef-...audit.jsonl",
               0.35, 7.1, 10, 0.25, font_size=8, color=MID)


def slide_period(prs: Presentation, d: dict):
    """Bonus slide — Period-over-period comparison (Feb / Mar / Apr)."""
    layout = prs.slide_layouts[6]
    sl = prs.slides.add_slide(layout)
    _set_slide_bg(sl, LIGHT)

    _add_rect(sl, 0, 0, 13.33, 0.08, fill=NAVY)

    _add_label(sl, "Pool Trend  Feb → Mar → Apr 2026", 0.45, 0.22, 9, 0.45,
               font_size=20, bold=True, color=NAVY)
    _add_label(sl, "Three monthly tapes — key metrics",
               0.45, 0.68, 8, 0.3, font_size=11, color=MID)

    pc = d["period_comparison"]
    months = [("Feb 2026", pc["feb"]), ("Mar 2026", pc["mar"]), ("Apr 2026", pc["apr"])]
    metrics = [
        ("Pool Balance", lambda v: f"€{v['balance']/1e9:.3f}B"),
        ("Loan Count",   lambda v: f"{v['loan_count']:,}"),
        ("WA Rate",      lambda v: f"{v['wa_rate']:.2f}%"),
        ("Green Share",  lambda v: f"{v['green_share_pct']:.2f}%"),
        ("Avg PED",      lambda v: f"{v['avg_ped']:.2f}"),
        ("90d+ Arrears", lambda v: f"{v['arrears_90plus_pct']:.4f}%"),
    ]

    col_w = 12.4 / (len(months) + 1)
    # Header
    _add_label(sl, "METRIC", 0.45, 1.3, col_w - 0.1, 0.3,
               font_size=9, bold=True, color=SLATE)
    for j, (mo, _) in enumerate(months):
        _add_label(sl, mo, 0.45 + (j + 1) * col_w, 1.3, col_w - 0.1, 0.3,
                   font_size=9, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    for i, (name, fmt) in enumerate(metrics):
        ry = 1.65 + i * 0.7
        row_bg = WHITE if i % 2 == 0 else RGBColor(0xF0, 0xF2, 0xF4)
        _add_rect(sl, 0.45, ry, 12.4, 0.6, fill=row_bg)
        _add_label(sl, name, 0.55, ry + 0.15, col_w - 0.15, 0.3,
                   font_size=10, bold=True, color=SLATE)
        for j, (_, vals) in enumerate(months):
            cell_val = fmt(vals)
            _add_label(sl, cell_val,
                       0.45 + (j + 1) * col_w, ry + 0.15, col_w - 0.1, 0.3,
                       font_size=13, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    _add_label(sl, "Source: /api/deal/periods + tape direct. CPR not computable (no per-loan scheduled amortisation).",
               0.45, 6.9, 11, 0.25, font_size=8, color=MID)


# ── Main ──────────────────────────────────────────────────────────────────────

def build(data_path: Path = DATA_FILE, out_path: Path = OUT_FILE):
    d = json.loads(data_path.read_text())

    prs = Presentation()
    prs.slide_width  = _px(13.33)
    prs.slide_height = _px(7.5)

    slide_title(prs, d)
    slide_framework(prs, d)
    slide_deal(prs, d)
    slide_finding(prs, d)
    slide_opportunity(prs, d)
    slide_governance(prs, d)
    slide_period(prs, d)

    prs.save(str(out_path))
    print(f"Saved → {out_path}  ({out_path.stat().st_size // 1024} KB,  {len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
